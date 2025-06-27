# Arquivo: main.py (Versão com OpenAI/ChatGPT)

import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
import json

# --- Importação das Ferramentas e Bibliotecas ---
# Importação para OpenAI
from openai import OpenAI
# Importação para Pesquisa na Web
from tavily import TavilyClient
# Importação para processamento de arquivos
from PIL import Image
import fitz # PyMuPDF para PDF
from docx import Document # para .docx
from pptx import Presentation # para .pptx

# --- CONFIGURAÇÃO INICIAL ---
app = Flask(__name__)
CORS(app)

# Pasta para salvar temporariamente os arquivos de upload
UPLOAD_FOLDER = 'temp_uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# --- CONFIGURAÇÃO DAS APIS ---
# Carrega as chaves das variáveis de ambiente para segurança
try:
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    tavily_api_key = os.environ.get('TAVILY_API_KEY')

    # Inicializa o cliente da OpenAI
    client = OpenAI(api_key=openai_api_key)

    if not openai_api_key or not tavily_api_key:
        print("AVISO: Uma ou mais chaves de API (OPENAI ou TAVILY) não foram encontradas nas variáveis de ambiente.")
        
except Exception as e:
    print(f"ERRO CRÍTICO DE CONFIGURAÇÃO DE API: {e}")

# --- DEFINIÇÃO DA FERRAMENTA DE PESQUISA WEB (A MESMA!) ---
def search_web(query: str):
    """
    Pesquisa na internet usando a API Tavily para obter informações em tempo real.
    Use esta ferramenta para responder perguntas sobre eventos atuais, notícias,
    ou quando precisar de informações que não estão no seu conhecimento interno.
    """
    try:
        tavily_client = TavilyClient(api_key=tavily_api_key)
        print(f"Chamando a ferramenta de pesquisa com a query: '{query}'")
        response = tavily_client.search(query=query, search_depth="advanced", max_results=5)
        return json.dumps(response.get('results', [])) # Retorna como JSON string para o modelo
    except Exception as e:
        print(f"Erro na ferramenta de pesquisa web: {e}")
        return f"Ocorreu um erro ao tentar pesquisar: {e}"

# --- DEFINIÇÃO DA FERRAMENTA PARA O MODELO DA OPENAI ---
# Este é o formato que a OpenAI precisa para "ver" a ferramenta
tools = [
    {
        "type": "function",
        "function": {
            "name": "search_web",
            "description": "Pesquisa na internet usando a API Tavily para obter informações em tempo real. Use para responder sobre eventos atuais ou informações que você não possui.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "O tópico ou pergunta a ser pesquisado na internet.",
                    },
                },
                "required": ["query"],
            },
        },
    }
]

# --- FUNÇÕES DE PROCESSAMENTO DE ARQUIVOS ---
# (Manter o código das funções de imagem, PDF, etc., é o mesmo)
def analisar_imagem_com_gemini(arquivo_stream):
    # Nota: Esta função ainda usa o Gemini. Você precisaria da API de visão da OpenAI para substituí-la.
    # Por agora, mantenha-a assim ou remova-a se não for usar.
    return "Função de análise de imagem temporariamente indisponível (usa API do Gemini)."

def extract_text_from_pdf(file_path):
    # ... mesmo código ...
    try:
        doc = fitz.open(file_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
        return text
    except Exception as e:
        return f"Erro ao processar PDF: {e}"

def extract_text_from_docx(file_path):
    # ... mesmo código ...
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Erro ao processar DOCX: {e}"

def extract_text_from_pptx(file_path):
    # ... mesmo código ...
    try:
        prs = Presentation(file_path)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text
    except Exception as e:
        return f"Erro ao processar PPTX: {e}"

# --- ROTAS DA API ---

@app.route('/chat', methods=['POST'])
def handle_chat():
    """
    Rota para conversas de texto, com capacidade de pesquisa na web, usando a API da OpenAI.
    """
    data = request.get_json()
    if not data or 'message' not in data:
        return jsonify({'erro': 'Mensagem não encontrada'}), 400
    
    user_message = data['message']
    
    try:
        # Envia a mensagem para o modelo da OpenAI com as ferramentas
        response = client.chat.completions.create(
            model="gpt-3.5-turbo", # Ou "gpt-4-turbo" para resultados melhores
            messages=[{"role": "user", "content": user_message}],
            tools=tools,
            tool_choice="auto", # Permite que o modelo decida se usa a ferramenta
        )

        # O primeiro objeto da resposta que vem da API
        response_message = response.choices[0].message
        
        # Verifica se o modelo quer chamar uma ferramenta
        if response_message.tool_calls:
            # Obtém a chamada de função
            tool_call = response_message.tool_calls[0]
            function_name = tool_call.function.name
            function_args = json.loads(tool_call.function.arguments)

            # Verifica qual função foi chamada
            if function_name == "search_web":
                query = function_args.get("query")
                
                # Executa a função e obtém o resultado
                function_response = search_web(query)

                # Envia a resposta da ferramenta de volta para o modelo
                second_response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "user", "content": user_message},
                        response_message, # Inclui a mensagem que pediu a chamada
                        {
                            "role": "tool",
                            "tool_call_id": tool_call.id,
                            "name": function_name,
                            "content": function_response, # O conteúdo da pesquisa
                        },
                    ],
                )
                return jsonify({'response': second_response.choices[0].message.content})
        
        # Se nenhuma ferramenta for chamada, retorna a resposta normal
        return jsonify({'response': response_message.content})

    except Exception as e:
        print(f"Erro na rota /chat: {e}")
        return jsonify({'erro': f"Ocorreu um erro no servidor: {e}"}), 500

@app.route('/reconhecer', methods=['POST'])
def upload_e_reconhecer():
    """
    Rota para upload e análise de arquivos.
    Nota: A análise de imagem ainda usa o Gemini. A análise de texto funcionará.
    """
    if 'file' not in request.files:
        return jsonify({'erro': 'Nenhum arquivo enviado'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'erro': 'Nome de arquivo vazio'}), 400

# Arquivo: main.py (Versão com OpenAI/ChatGPT)

import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
import json
import base64

# --- Importação das Ferramentas e Bibliotecas ---
# Importação para OpenAI
from openai import OpenAI
# Importação para Pesquisa na Web
from tavily import TavilyClient
# Importação para processamento de arquivos
from PIL import Image
import fitz # PyMuPDF para PDF
from docx import Document # para .docx
from pptx import Presentation # para .pptx

# --- CONFIGURAÇÃO INICIAL ---
app = Flask(__name__)
CORS(app)

# Pasta para salvar temporariamente os arquivos de upload
UPLOAD_FOLDER = 'temp_uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# --- CONFIGURAÇÃO DAS APIS ---
# Carrega as chaves das variáveis de ambiente para segurança
try:
    openai_api_key = os.environ.get('OPENAI_API_KEY')
    tavily_api_key = os.environ.get('TAVILY_API_KEY')

    # Inicializa o cliente da OpenAI
    client = OpenAI(api_key=openai_api_key)
    
    # Inicializa o cliente da Tavily
    tavily_client = TavilyClient(api_key=tavily_api_key)

    if not openai_api_key or not tavily_api_key:
        print("AVISO: Uma ou mais chaves de API (OPENAI ou TAVILY) não foram encontradas nas variáveis de ambiente.")
        
except Exception as e:
    print(f"ERRO CRÍTICO DE CONFIGURAÇÃO DE API: {e}")

# --- DEFINIÇÃO DA FERRAMENTA DE PESQUISA WEB ---
def search_web(query: str):
    """
    Pesquisa na internet usando a API Tavily para obter informações em tempo real.
    Use esta ferramenta para responder perguntas sobre eventos atuais, notícias,
    ou quando precisar de informações que não estão no seu conhecimento interno.
    """
    try:
        print(f"Chamando a ferramenta de pesquisa com a query: '{query}'")
        response = tavily_client.search(query=query, search_depth="advanced", max_results=5)
        # Retorna os resultados como uma string JSON para o modelo
        return json.dumps(response.get('results', []))
    except Exception as e:
        print(f"Erro na ferramenta de pesquisa web: {e}")
        return f"Ocorreu um erro ao tentar pesquisar: {e}"

# --- DEFINIÇÃO DA FERRAMENTA PARA O MODELO DA OPENAI ---
# Este é o formato que a OpenAI precisa para "ver" a ferramenta
tools = [
    {
        "type": "function",
        "function": {
            "name": "search_web",
            "description": "Pesquisa na internet usando a API Tavily para obter informações em tempo real. Use para responder sobre eventos atuais ou informações que você não possui.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "O tópico ou pergunta a ser pesquisado na internet.",
                    },
                },
                "required": ["query"],
            },
        },
    }
]

# --- FUNÇÕES DE PROCESSAMENTO DE ARQUIVOS ---
def analisar_imagem_com_openai(arquivo_stream):
    """Envia uma imagem para a API de Visão da OpenAI e pede uma descrição."""
    try:
        # Codifica o stream do arquivo em base64
        base64_image = base64.b64encode(arquivo_stream.read()).decode('utf-8')
        
        # Faz a chamada para a API de Visão do GPT-4
        response = client.chat.completions.create(
            model="gpt-4o", # gpt-4o ou gpt-4-vision-preview
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": "Descreva detalhadamente esta imagem. Se for um componente industrial, explique o que é, sua função e possíveis falhas. Se não, apenas descreva o que você vê."},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{base64_image}",
                                "detail": "high"
                            },
                        },
                    ],
                }
            ],
            max_tokens=500,
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Ocorreu um erro ao tentar analisar a imagem com a OpenAI: {e}"

def extract_text_from_pdf(file_path):
    """Extrai texto de um arquivo PDF."""
    try:
        doc = fitz.open(file_path)
        text = "".join(page.get_text() for page in doc)
        doc.close()
        return text
    except Exception as e:
        return f"Erro ao processar PDF: {e}"

def extract_text_from_docx(file_path):
    """Extrai texto de um arquivo .docx."""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"Erro ao processar DOCX: {e}"

def extract_text_from_pptx(file_path):
    """Extrai texto de um arquivo .pptx."""
    try:
        prs = Presentation(file_path)
        text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        return text
    except Exception as e:
        return f"Erro ao processar PPTX: {e}"

# --- ROTAS DA API ---

@app.route('/chat', methods=['POST'])
def handle_chat():
    """
    Rota para conversas de texto, com capacidade de pesquisa na web, usando a API da OpenAI.
    """
    data = request.get_json()
    if not data or 'message' not in data:
        return jsonify({'erro': 'Mensagem não encontrada'}), 400
    
    user_message = data['message']
    
    try:
        print(f"Requisição POST recebida na rota /chat")
        
        # Envia a mensagem para o modelo da OpenAI com as ferramentas
        response = client.chat.completions.create(
            model="gpt-3.5-turbo", # Use "gpt-4o" se preferir um modelo mais avançado
            messages=[{"role": "user", "content": user_message}],
            tools=tools,
            tool_choice="auto", # Permite que o modelo decida se usa a ferramenta
        )

        response_message = response.choices[0].message
        
        # Verifica se o modelo quer chamar uma ferramenta
        if response_message.tool_calls:
            tool_call = response_message.tool_calls[0]
            function_name = tool_call.function.name
            function_args = json.loads(tool_call.function.arguments)

            if function_name == "search_web":
                query = function_args.get("query")
                
                # Executa a função e obtém o resultado
                function_response = search_web(query)

                # Envia a resposta da ferramenta de volta para o modelo
                second_response = client.chat.completions.create(
                    model="gpt-3.5-turbo", # Use o mesmo modelo
                    messages=[
                        {"role": "user", "content": user_message},
                        response_message, # Inclui a mensagem que pediu a chamada
                        {
                            "role": "tool",
                            "tool_call_id": tool_call.id,
                            "name": function_name,
                            "content": function_response, # O conteúdo da pesquisa
                        },
                    ],
                )
                return jsonify({'response': second_response.choices[0].message.content})
        
        # Se nenhuma ferramenta for chamada, retorna a resposta normal
        return jsonify({'response': response_message.content})

    except Exception as e:
        print(f"Erro na rota /chat: {e}")
        return jsonify({'erro': f"Ocorreu um erro no servidor: {e}"}), 500

@app.route('/reconhecer', methods=['POST'])
def upload_e_reconhecer():
    """Rota para upload e análise de arquivos (imagens, PDF, DOCX, PPTX)."""
    if 'file' not in request.files:
        return jsonify({'erro': 'Nenhum arquivo enviado'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'erro': 'Nome de arquivo vazio'}), 400

    filename = secure_filename(file.filename)
    file_extension = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    conteudo_reconhecido = ""

    if file_extension in ['png', 'jpg', 'jpeg', 'webp']:
        # Chama a função de análise de imagem do OpenAI
        conteudo_reconhecido = analisar_imagem_com_openai(file.stream)
    elif file_extension in ['pdf', 'docx', 'pptx']:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        try:
            if file_extension == 'pdf':
                conteudo_reconhecido = extract_text_from_pdf(file_path)
            elif file_extension == 'docx':
                conteudo_reconhecido = extract_text_from_docx(file_path)
            elif file_extension == 'pptx':
                conteudo_reconhecido = extract_text_from_pptx(file_path)
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)
    else:
        conteudo_reconhecido = f"Tipo de arquivo ('{file_extension}') não suportado."
    
    return jsonify({'conteudo_extraido': conteudo_reconhecido})

@app.route('/')
def index():
    """Rota principal para verificar se a API está online."""
    return "API da AEMI v4.0 com OpenAI e análise de arquivos está no ar!"

# --- INICIALIZAÇÃO DO SERVIDOR ---
if __name__ == '__main__':
    port = int(os.environ.get('PORT', 8000))
    app.run(host='0.0.0.0', port=port)

